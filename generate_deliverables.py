"""
Generateur automatique des livrables du projet :
  1. Rapport professionnel  -> reports/Rapport_ROI_Marketing.docx
  2. Presentation PowerPoint -> reports/Presentation_ROI_Marketing.pptx

Usage :
    python generate_deliverables.py
"""

import io
from pathlib import Path

import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGB
from pptx.util import Emu, Inches as PInches, Pt as PPt

PROJECT_ROOT = Path(__file__).resolve().parent
FIG = PROJECT_ROOT / "reports" / "figures"
REPORTS = PROJECT_ROOT / "reports"

# ─── Vraies données ──────────────────────────────────────────────────────────
eval_df  = pd.read_csv(REPORTS / "evaluation_results.csv")
train_df = pd.read_csv(PROJECT_ROOT / "models" / "training_results.csv")
interp_df = pd.read_csv(REPORTS / "interpretability_results.csv")

# ─── Palette ─────────────────────────────────────────────────────────────────
C_BLUE_D  = RGBColor(0x1e, 0x3a, 0x5f)
C_BLUE_M  = RGBColor(0x25, 0x63, 0xeb)
C_BLUE_L  = RGBColor(0xdb, 0xea, 0xfe)
C_GREEN   = RGBColor(0x10, 0xb9, 0x81)
C_AMBER   = RGBColor(0xf5, 0x9e, 0x0b)
C_RED     = RGBColor(0xe7, 0x4c, 0x3c)
C_WHITE   = RGBColor(0xff, 0xff, 0xff)
C_GRAY    = RGBColor(0xf1, 0xf5, 0xf9)
C_DARK    = RGBColor(0x1e, 0x29, 0x3b)
C_MUTED   = RGBColor(0x64, 0x74, 0x8b)


# =============================================================================
# HELPERS DOCX
# =============================================================================

def set_cell_bg(cell, rgb: RGBColor):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    hex_color = f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color)
    tcPr.append(shd)


def set_cell_border(cell, border_size=4, color="2563eb"):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ("top", "bottom", "left", "right"):
        tcBdr = OxmlElement(f"w:{edge}")
        tcBdr.set(qn("w:val"), "single")
        tcBdr.set(qn("w:sz"), str(border_size))
        tcBdr.set(qn("w:color"), color)
        tcBdr.set(qn("w:space"), "0")
        if tcPr.find(qn(f"w:{edge}")) is None:
            bdr = OxmlElement("w:tcBdr")
            bdr.append(tcBdr)
            tcPr.append(bdr)


def para_style(para, bold=False, italic=False, size_pt=11,
               color: RGBColor = None, align=None, space_after=0):
    run = para.runs[0] if para.runs else para.add_run("")
    run.bold = bold
    run.italic = italic
    run.font.size = Pt(size_pt)
    if color:
        run.font.color.rgb = color
    if align:
        para.alignment = align
    para.paragraph_format.space_after = Pt(space_after)
    return para


def add_heading(doc, text, level=1, color=None):
    h = doc.add_heading(text, level=level)
    h.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for run in h.runs:
        run.font.color.rgb = color or C_BLUE_D
        run.font.name = "Calibri"
    return h


def add_image_safe(doc, path, width=Inches(5.5)):
    p = Path(path)
    if p.exists():
        doc.add_picture(str(p), width=width)
        last = doc.paragraphs[-1]
        last.alignment = WD_ALIGN_PARAGRAPH.CENTER
    else:
        doc.add_paragraph(f"[Figure non disponible : {p.name}]")


def add_metric_table(doc, headers, rows, header_color=None):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = "Table Grid"
    hdr_row = table.rows[0]
    for i, h in enumerate(headers):
        cell = hdr_row.cells[i]
        cell.text = h
        set_cell_bg(cell, header_color or C_BLUE_D)
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.runs[0]
        run.bold = True
        run.font.color.rgb = C_WHITE
        run.font.size = Pt(10)
        run.font.name = "Calibri"

    for ri, row_data in enumerate(rows):
        row = table.rows[ri + 1]
        bg = C_GRAY if ri % 2 == 0 else C_WHITE
        for ci, val in enumerate(row_data):
            cell = row.cells[ci]
            cell.text = str(val)
            set_cell_bg(cell, bg)
            p = cell.paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.runs[0]
            run.font.size = Pt(10)
            run.font.name = "Calibri"
            # Surligner la meilleure ligne (XGBoost = ri 0)
            if ri == 0:
                run.font.color.rgb = C_GREEN
                run.bold = True

    doc.add_paragraph("")
    return table


def add_colored_box(doc, text, bg: RGBColor, text_color: RGBColor = None):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    p.paragraph_format.space_before = Pt(4)
    p.paragraph_format.space_after = Pt(4)
    p.paragraph_format.left_indent = Inches(0.15)
    p.paragraph_format.right_indent = Inches(0.15)
    run = p.add_run(text)
    run.font.size = Pt(10.5)
    run.font.name = "Calibri"
    run.font.color.rgb = text_color or C_DARK
    pPr = p._p.get_or_add_pPr()
    shd = OxmlElement("w:shd")
    hex_c = f"{bg[0]:02X}{bg[1]:02X}{bg[2]:02X}"
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_c)
    pPr.append(shd)


# =============================================================================
# RAPPORT WORD
# =============================================================================

def build_report():
    doc = Document()
    section = doc.sections[0]
    section.page_width  = Inches(8.27)
    section.page_height = Inches(11.69)
    section.left_margin  = Inches(1.1)
    section.right_margin = Inches(1.1)
    section.top_margin   = Inches(1.1)
    section.bottom_margin = Inches(0.9)

    # ── Style corps de texte ──────────────────────────────────────────────────
    normal_style = doc.styles["Normal"]
    normal_style.font.name = "Calibri"
    normal_style.font.size = Pt(11)

    # ─────────────────────────────────────────────────────────────────────────
    # PAGE DE GARDE
    # ─────────────────────────────────────────────────────────────────────────
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(60)
    run = p.add_run("EFREI Paris  |  DATA Engineering & AI")
    run.bold = True
    run.font.size = Pt(13)
    run.font.color.rgb = C_MUTED
    run.font.name = "Calibri"

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Master 1 — Projet Data Science")
    run.font.size = Pt(12)
    run.font.color.rgb = C_MUTED
    run.font.name = "Calibri"

    doc.add_paragraph("")
    doc.add_paragraph("")

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Système Intelligent Multi-Modèles")
    run.bold = True
    run.font.size = Pt(26)
    run.font.color.rgb = C_BLUE_D
    run.font.name = "Calibri"

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Optimisation du Retour sur Investissement Marketing")
    run.bold = True
    run.font.size = Pt(20)
    run.font.color.rgb = C_BLUE_M
    run.font.name = "Calibri"

    doc.add_paragraph("")

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("─" * 45)
    run.font.color.rgb = C_BLUE_L
    run.font.name = "Calibri"

    doc.add_paragraph("")
    doc.add_paragraph("")

    for label, value in [
        ("Auteur", "Salah Eddine Zbiri"),
        ("Encadrant", "Sarah Malaeb"),
        ("Promotion", "M1 Data Engineering & AI — 2025/2026"),
        ("Date de rendu", "Juin 2026"),
    ]:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r1 = p.add_run(f"{label} : ")
        r1.font.color.rgb = C_MUTED
        r1.font.size = Pt(12)
        r1.font.name = "Calibri"
        r2 = p.add_run(value)
        r2.bold = True
        r2.font.size = Pt(12)
        r2.font.color.rgb = C_DARK
        r2.font.name = "Calibri"

    doc.add_page_break()

    # ─────────────────────────────────────────────────────────────────────────
    # RÉSUMÉ EXÉCUTIF
    # ─────────────────────────────────────────────────────────────────────────
    add_heading(doc, "Résumé Exécutif", level=1)
    add_colored_box(doc,
        "Ce projet présente la conception d'un système intelligent complet d'aide à la décision "
        "pour l'optimisation du ROI marketing. À partir du dataset Kaggle Dummy Advertising & Sales "
        "(~4 500 campagnes multicanal : TV, Radio, Social Media, Influencer), quatre modèles supervisés "
        "ont été entraînés, comparés et interprétés. Le modèle XGBoost retenu atteint un R² de 0.9985 "
        "et une MAPE de 1.84% sur le test set, avec zéro overfitting. L'analyse SHAP révèle que le "
        "budget TV explique à lui seul 95.69% des ventes, tandis que le type d'influenceur n'a qu'un "
        "impact de 0.13% — une information stratégique majeure pour le CMO. La solution inclut un "
        "dashboard Streamlit interactif (5 pages) et une API REST FastAPI déployable.",
        bg=C_BLUE_L, text_color=C_BLUE_D
    )
    doc.add_paragraph("")

    # ─────────────────────────────────────────────────────────────────────────
    # TABLE DES MATIÈRES (manuelle)
    # ─────────────────────────────────────────────────────────────────────────
    add_heading(doc, "Table des matières", level=1)
    toc_items = [
        ("1.", "Introduction et contexte métier", "3"),
        ("2.", "Dataset et Analyse Exploratoire (EDA)", "4"),
        ("3.", "Pipeline de Preprocessing", "6"),
        ("4.", "Modélisation multi-algorithmes", "8"),
        ("5.", "Évaluation et comparaison des modèles", "10"),
        ("6.", "Interprétabilité des prédictions", "13"),
        ("7.", "Dashboard décisionnel (Streamlit)", "15"),
        ("8.", "API REST (FastAPI)", "16"),
        ("9.", "Conclusion et recommandations", "17"),
    ]
    for num, title, page in toc_items:
        p = doc.add_paragraph()
        p.paragraph_format.tab_stops.add_tab_stop(Inches(5.5))
        r1 = p.add_run(f"{num}  {title}")
        r1.font.size = Pt(11)
        r1.font.name = "Calibri"
        r2 = p.add_run(f"\t{page}")
        r2.font.size = Pt(11)
        r2.font.color.rgb = C_MUTED
        r2.font.name = "Calibri"

    doc.add_page_break()

    # ─────────────────────────────────────────────────────────────────────────
    # 1. INTRODUCTION
    # ─────────────────────────────────────────────────────────────────────────
    add_heading(doc, "1. Introduction et contexte métier", level=1)
    add_heading(doc, "1.1 Problématique", level=2, color=C_BLUE_M)
    doc.add_paragraph(
        "Dans un environnement marketing multicanal hautement concurrentiel, les entreprises investissent "
        "des millions de dollars en budgets publicitaires (TV, Radio, Social Media, Influencer) sans "
        "disposer d'outils objectifs pour estimer le retour sur investissement (ROI) avant l'exécution "
        "des campagnes. La question centrale de ce projet est la suivante :\n"
        "Comment prédire le volume de ventes généré par une combinaison donnée de budgets marketing, "
        "et quels canaux maximisent réellement le ROI ?"
    )
    add_heading(doc, "1.2 Objectifs pédagogiques", level=2, color=C_BLUE_M)
    for item in [
        "Implémenter un pipeline Data Science complet : EDA → Preprocessing → Modélisation → Évaluation → Interprétabilité",
        "Comparer au minimum 4 modèles supervisés incluant un modèle Deep Learning (MLP)",
        "Identifier le meilleur compromis performance / interprétabilité / coût computationnel",
        "Concevoir un dashboard décisionnel orienté utilisateur métier (CMO)",
        "Industrialiser la solution via une API REST déployable (FastAPI)",
    ]:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(item)
        run.font.size = Pt(11)
        run.font.name = "Calibri"

    add_heading(doc, "1.3 Compétences RNCP40875 — Bloc 2 ciblées", level=2, color=C_BLUE_M)
    add_metric_table(doc,
        ["Compétence", "Description"],
        [
            ["C3.1", "Préparation et transformation des données (pipeline sklearn anti data leakage)"],
            ["C3.2", "Communication infographique (dashboard 5 pages, graphiques interactifs)"],
            ["C3.3", "Analyse exploratoire documentée (EDA avec insights métier)"],
            ["C4.1", "Stratégie d'intégration IA (API REST + dashboard orienté CMO)"],
            ["C4.2", "Développement de modèles prédictifs (4 modèles, GridSearchCV)"],
            ["C4.3", "Évaluation et comparaison multi-modèles (métriques + overfitting)"],
        ]
    )

    doc.add_page_break()

    # ─────────────────────────────────────────────────────────────────────────
    # 2. EDA
    # ─────────────────────────────────────────────────────────────────────────
    add_heading(doc, "2. Dataset et Analyse Exploratoire (EDA)", level=1)
    add_heading(doc, "2.1 Présentation du dataset", level=2, color=C_BLUE_M)
    doc.add_paragraph(
        "Le dataset utilisé est le Dummy Advertising and Sales Data disponible sur Kaggle "
        "(harrimansaragih/dummy-advertising-and-sales-data). Il simule un environnement marketing "
        "réaliste avec des données synthétiques cohérentes."
    )
    add_metric_table(doc,
        ["Variable", "Type", "Description", "Plage"],
        [
            ["TV", "float", "Budget télévision (millions $)", "0 — ~296 M$"],
            ["Radio", "float", "Budget radio (millions $)", "0 — ~50 M$"],
            ["Social Media", "float", "Budget social media (millions $)", "0 — ~175 M$"],
            ["Influencer", "string", "Type : Nano / Micro / Macro / Mega", "4 catégories"],
            ["Sales (cible)", "float", "Ventes générées (millions $)", "~0 — ~483 M$"],
        ]
    )

    add_heading(doc, "2.2 Valeurs manquantes et qualité", level=2, color=C_BLUE_M)
    doc.add_paragraph(
        "L'analyse des valeurs manquantes révèle un taux inférieur à 5% uniquement sur la variable TV. "
        "Aucun doublon n'a été détecté. Les distributions sont légèrement asymétriques à droite pour TV "
        "(skewness positif), ce qui justifie l'imputation par la médiane plutôt que par la moyenne."
    )
    add_image_safe(doc, FIG / "01_missing_values.png", Inches(4.5))
    doc.add_paragraph("Figure 1 — Valeurs manquantes par colonne").alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_heading(doc, "2.3 Distribution de la variable cible", level=2, color=C_BLUE_M)
    doc.add_paragraph(
        "La variable cible Sales présente une distribution bi-modale légère avec une concentration "
        "entre 0 et 100 M$ et quelques valeurs élevées jusqu'à 483 M$. La médiane (≈ 12 M$) est "
        "proche de la moyenne, indiquant une distribution relativement symétrique sans distorsion majeure."
    )
    add_image_safe(doc, FIG / "02_sales_distribution.png", Inches(5.0))
    doc.add_paragraph("Figure 2 — Distribution de Sales (histogramme + boxplot)").alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_heading(doc, "2.4 Relations bivariées et corrélations", level=2, color=C_BLUE_M)
    add_colored_box(doc,
        "RÉSULTAT CLEF — Corrélation TV-Sales = 0.999 (quasi-parfaite). "
        "Ce signal est exceptionnellement fort et suggère une relation quasi-linéaire entre le budget TV "
        "et les ventes. Radio présente une corrélation de 0.87, Social Media de 0.67. "
        "Le type d'Influencer a un impact marginal non significatif.",
        bg=RGBColor(0xd1, 0xfa, 0xe5), text_color=RGBColor(0x06, 0x5f, 0x46)
    )
    add_image_safe(doc, FIG / "05_budgets_vs_sales.png", Inches(5.5))
    doc.add_paragraph("Figure 3 — Scatter plots Budget vs Sales (avec droite de régression)").alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_image_safe(doc, FIG / "07_correlation_matrix.png", Inches(4.0))
    doc.add_paragraph("Figure 4 — Matrice de corrélation (Pearson)").alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_heading(doc, "2.5 Impact du type d'Influencer", level=2, color=C_BLUE_M)
    doc.add_paragraph(
        "L'analyse des ventes moyennes par type d'influenceur (Nano, Micro, Macro, Mega) ne révèle "
        "aucune différence statistiquement significative. Ce résultat, confirmé ultérieurement par SHAP "
        "(0.13% d'importance), constitue un insight business majeur : les dépenses en Mega influenceurs "
        "ne se justifient pas par leur impact sur les ventes."
    )
    add_image_safe(doc, FIG / "06_influencer_vs_sales.png", Inches(5.0))
    doc.add_paragraph("Figure 5 — Distribution Sales par type d'Influencer").alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_page_break()

    # ─────────────────────────────────────────────────────────────────────────
    # 3. PREPROCESSING
    # ─────────────────────────────────────────────────────────────────────────
    add_heading(doc, "3. Pipeline de Preprocessing", level=1)
    add_heading(doc, "3.1 Principe d'anti data leakage", level=2, color=C_BLUE_M)
    add_colored_box(doc,
        "PRINCIPE FONDAMENTAL : le preprocessor (imputation, scaling, encodage) est FIT "
        "UNIQUEMENT sur le train set et appliqué en transform sur le test set. "
        "Les pipelines sklearn garantissent ce principe même pendant la cross-validation "
        "(le preprocessing est refit à chaque fold sur les seules données train du fold).",
        bg=RGBColor(0xfe, 0xf3, 0xc7), text_color=RGBColor(0x78, 0x35, 0x00)
    )
    add_heading(doc, "3.2 Architecture du pipeline", level=2, color=C_BLUE_M)
    add_metric_table(doc,
        ["Étape", "Variables", "Transformation", "Justification"],
        [
            ["1 — Imputation", "TV, Radio, Social Media", "SimpleImputer(strategy='median')", "Robuste aux outliers (mieux que mean)"],
            ["2 — Standardisation", "TV, Radio, Social Media", "StandardScaler", "Requis pour Régression Linéaire et MLP"],
            ["3 — Encodage", "Influencer", "OneHotEncoder (4 catégories)", "Variable nominale, pas d'ordre ordinal strict"],
            ["4 — Split", "Tout le dataset", "train_test_split(80/20, seed=42)", "Reproductibilité garantie"],
        ]
    )
    add_heading(doc, "3.3 Résultats du preprocessing", level=2, color=C_BLUE_M)
    doc.add_paragraph(
        "Après transformation : 0 valeur manquante, features numériques avec mean ≈ 0 et std ≈ 1 "
        "sur le train (vérification post-transformation). Les 4 catégories Influencer sont encodées "
        "en 4 colonnes binaires. Le handle_unknown='ignore' sur l'OneHotEncoder garantit la robustesse "
        "face à des catégories inconnues en production (dashboard, API)."
    )

    doc.add_page_break()

    # ─────────────────────────────────────────────────────────────────────────
    # 4. MODÉLISATION
    # ─────────────────────────────────────────────────────────────────────────
    add_heading(doc, "4. Modélisation multi-algorithmes", level=1)
    add_heading(doc, "4.1 Stratégie de modélisation", level=2, color=C_BLUE_M)
    doc.add_paragraph(
        "Conformément aux exigences du sujet, quatre modèles sont entraînés et comparés. "
        "Chaque modèle est encapsulé dans un sklearn Pipeline complet (preprocessor + estimator) "
        "pour garantir l'absence de data leakage même pendant la GridSearchCV. "
        "La cross-validation KFold 5-fold est utilisée (shuffle, random_state=42). "
        "La métrique de sélection principale est le RMSE (Root Mean Squared Error)."
    )
    add_heading(doc, "4.2 Les 4 modèles implémentés", level=2, color=C_BLUE_M)
    models_info = [
        ("Régression Linéaire (Baseline)",
         "Modèle de référence. Aucun hyperparamètre à tuner. Sert de point de comparaison pour "
         "mesurer le gain apporté par les modèles complexes. Très interprétable.",
         "Aucune (modèle déterminé entièrement)"),
        ("Random Forest",
         "Ensemble de 200 arbres de décision entraînés en parallèle (bagging). Capture les non-linéarités "
         "et interactions sans feature engineering. Robuste aux outliers.",
         "n_estimators=200, max_depth=10, min_samples_split=2"),
        ("XGBoost (Gradient Boosting)",
         "Algorithme de boosting itératif. Étatde l'art sur les données tabulaires. "
         "Gradient descent sur une fonction objectif avec régularisation intégrée. "
         "Meilleur compromis performance/vitesse.",
         "n_estimators=200, max_depth=3, learning_rate=0.1"),
        ("MLP — Deep Learning (sklearn)",
         "Réseau de neurones multicouches (Multi-Layer Perceptron) avec backpropagation. "
         "Architecture : input → 2 couches cachées (50, 25) → output. Early stopping activé "
         "(anti-overfitting), régularisation L2 (alpha=0.001).",
         "hidden_layer_sizes=(50,25), alpha=0.001, lr_init=0.01, early_stopping=True"),
    ]
    for i, (name, desc, params) in enumerate(models_info):
        add_heading(doc, f"4.2.{i+1} {name}", level=3, color=C_DARK)
        doc.add_paragraph(desc)
        p = doc.add_paragraph()
        r = p.add_run(f"Meilleurs hyperparamètres (GridSearchCV) : {params}")
        r.italic = True
        r.font.color.rgb = C_MUTED
        r.font.size = Pt(10)
        r.font.name = "Calibri"
        doc.add_paragraph("")

    add_heading(doc, "4.3 Résultats de la cross-validation (5-fold, train set)", level=2, color=C_BLUE_M)
    cv_rows = []
    for _, row in train_df.sort_values("CV_RMSE").iterrows():
        cv_rows.append([row["Model"], f"{row['CV_RMSE']:.4f}", f"{row['CV_R2']:.4f}", f"{row['Time_sec']:.1f}s"])
    add_metric_table(doc,
        ["Modèle", "CV RMSE", "CV R²", "Temps"],
        cv_rows
    )
    doc.add_paragraph(
        "Note : la Régression Linéaire ne fait pas l'objet d'une GridSearchCV (pas d'hyperparamètre). "
        "Les résultats CV sont obtenus via cross_validate() avec le même KFold 5-fold."
    ).runs[0].font.italic = True

    doc.add_page_break()

    # ─────────────────────────────────────────────────────────────────────────
    # 5. ÉVALUATION
    # ─────────────────────────────────────────────────────────────────────────
    add_heading(doc, "5. Évaluation et comparaison des modèles", level=1)
    add_heading(doc, "5.1 Métriques utilisées", level=2, color=C_BLUE_M)
    add_metric_table(doc,
        ["Métrique", "Formule", "Interprétation"],
        [
            ["MAE", "|y - ŷ| moyen", "Erreur absolue moyenne (en M$)"],
            ["RMSE", "√(MSE)", "Penalise les grandes erreurs — sensible aux outliers"],
            ["R²", "1 - SS_res/SS_tot", "Part de variance expliquée (1 = parfait)"],
            ["MAPE (%)", "|y-ŷ|/|y| × 100", "Erreur relative interprétable par le métier"],
        ]
    )

    add_heading(doc, "5.2 Résultats sur le test set", level=2, color=C_BLUE_M)
    eval_rows = []
    for _, row in eval_df.sort_values("RMSE_test").iterrows():
        eval_rows.append([
            row["Model"],
            f"{row['MAE_test']:.3f}", f"{row['RMSE_test']:.3f}",
            f"{row['R2_test']:.4f}", f"{row['MAPE_test_%']:.2f}%",
            f"{row['Gap_R2']:+.4f}", row["Verdict"],
        ])
    add_metric_table(doc,
        ["Modèle", "MAE", "RMSE", "R²", "MAPE", "Gap R²", "Verdict"],
        eval_rows
    )
    add_colored_box(doc,
        "MODÈLE RETENU : XGBoost — RMSE = 3.538, R² = 0.9985, MAPE = 1.84%. "
        "Gap R² = 0.000 → aucun overfitting. Le modèle prédit les ventes avec une erreur "
        "moyenne de 1.84%, ce qui offre une précision suffisante pour la planification budgétaire.",
        bg=RGBColor(0xd1, 0xfa, 0xe5), text_color=RGBColor(0x06, 0x5f, 0x46)
    )

    add_heading(doc, "5.3 Visualisations", level=2, color=C_BLUE_M)
    add_image_safe(doc, FIG / "04_predictions_vs_actual.png", Inches(5.5))
    doc.add_paragraph("Figure 6 — Prédictions vs Ventes réelles (4 modèles, test set)").alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_image_safe(doc, FIG / "07_metrics_comparison.png", Inches(5.5))
    doc.add_paragraph("Figure 7 — Comparaison des métriques (RMSE, R², MAPE)").alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_image_safe(doc, FIG / "05_residuals_distribution.png", Inches(5.5))
    doc.add_paragraph("Figure 8 — Distribution des résidus (4 modèles)").alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_heading(doc, "5.4 Analyse des résidus", level=2, color=C_BLUE_M)
    doc.add_paragraph(
        "L'analyse des résidus (y_réel − ŷ) confirme la qualité des prédictions XGBoost : "
        "les résidus sont centrés sur 0, symétriquement distribués et sans pattern structuré "
        "dans le graphe résidus vs prédictions. Cela indique que le modèle ne souffre "
        "d'aucun biais systématique, ni d'hétéroscédasticité notable."
    )
    add_image_safe(doc, FIG / "06_residuals_vs_predictions.png", Inches(5.0))
    doc.add_paragraph("Figure 9 — Résidus vs Prédictions (4 modèles)").alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_heading(doc, "5.5 Analyse comparative — Pourquoi XGBoost ?", level=2, color=C_BLUE_M)
    doc.add_paragraph(
        "Régression Linéaire : le R² de 0.9960 est remarquablement élevé pour un modèle si simple, "
        "confirmant la nature quasi-linéaire de la relation TV→Sales. Elle constitue un excellent "
        "baseline et reste recommandée pour l'interprétation directe des coefficients.\n\n"
        "Random Forest : légèrement supérieur en RMSE test (3.701 vs 3.538 pour XGBoost) mais "
        "affiche un Gap R² de +0.0007 (micro-overfitting). Entraînement 2× plus lent que XGBoost.\n\n"
        "MLP (Deep Learning) : performances inférieures au baseline linéaire (RMSE 5.983 vs 5.879). "
        "Illustration pédagogique cruciale : le Deep Learning n'est pas toujours supérieur. "
        "Sur un dataset de cette taille (~4 500 lignes) avec des relations essentiellement linéaires, "
        "un réseau de neurones apporte peu de valeur ajoutée.\n\n"
        "XGBoost : meilleur RMSE (3.538), zéro overfitting, 3.6s d'entraînement (le plus rapide "
        "des modèles complexes). Gradient boosting avec régularisation intégrée, optimal pour "
        "les données tabulaires."
    )

    doc.add_page_break()

    # ─────────────────────────────────────────────────────────────────────────
    # 6. INTERPRÉTABILITÉ
    # ─────────────────────────────────────────────────────────────────────────
    add_heading(doc, "6. Interprétabilité des prédictions", level=1)
    doc.add_paragraph(
        "Trois méthodes complémentaires sont appliquées sur le modèle final XGBoost "
        "pour expliquer ses décisions à différents niveaux de granularité."
    )

    add_heading(doc, "6.1 Feature Importance native (XGBoost — Gain)", level=2, color=C_BLUE_M)
    doc.add_paragraph(
        "Calculée gratuitement pendant l'entraînement. Mesure la réduction moyenne d'impureté "
        "(variance) lors de chaque split utilisant la variable. Rapide mais dépend de la structure "
        "des arbres et peut être biaisée en présence de variables corrélées."
    )
    add_image_safe(doc, FIG / "08_feature_importance_native.png", Inches(5.0))
    doc.add_paragraph("Figure 10 — Feature Importance native XGBoost (Gain)").alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_heading(doc, "6.2 Permutation Importance", level=2, color=C_BLUE_M)
    doc.add_paragraph(
        "Méthode agnostique au modèle appliquée sur le test set. Pour chaque variable, on mesure "
        "la chute de R² après permutation aléatoire de ses valeurs (10 répétitions). "
        "Plus la chute est grande, plus la variable est critique. Avantage : mesure l'impact réel "
        "sur la généralisation, pas seulement sur l'entraînement."
    )
    add_image_safe(doc, FIG / "09_permutation_importance.png", Inches(5.0))
    doc.add_paragraph("Figure 11 — Permutation Importance (chute R² — test set)").alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_heading(doc, "6.3 SHAP (SHapley Additive exPlanations)", level=2, color=C_BLUE_M)
    doc.add_paragraph(
        "Méthode la plus avancée, basée sur la théorie des jeux (valeurs de Shapley). "
        "TreeExplainer (optimisé pour XGBoost) calcule la contribution de chaque variable "
        "à chaque prédiction individuelle. Permet à la fois une vue globale (quelle variable "
        "influence le plus en moyenne) et une vue locale (pourquoi cette prédiction précise)."
    )
    add_image_safe(doc, FIG / "10_shap_importance_bar.png", Inches(5.0))
    doc.add_paragraph("Figure 12 — SHAP Importance globale (mean |SHAP value|)").alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_image_safe(doc, FIG / "11_shap_beeswarm.png", Inches(5.0))
    doc.add_paragraph("Figure 13 — SHAP Beeswarm (impact + direction)").alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_heading(doc, "6.4 Synthèse des 3 méthodes", level=2, color=C_BLUE_M)
    interp_rows = []
    for _, row in interp_df.iterrows():
        interp_rows.append([
            row["feature"],
            f"{row['Native_Importance_norm']:.2f}%",
            f"{row['Permutation_R2_loss_norm']:.2f}%",
            f"{row['SHAP_mean_abs_norm']:.2f}%",
        ])
    add_metric_table(doc,
        ["Variable", "Native (%)", "Permutation (%)", "SHAP (%)"],
        interp_rows
    )
    add_colored_box(doc,
        "CONVERGENCE DES 3 MÉTHODES : TV domine avec 95-99% d'importance selon la méthode. "
        "Radio en second plan (1-4%). Social Media et Influencer négligeables (<1%). "
        "Cette convergence est un signal de robustesse de l'analyse — les 3 méthodes "
        "indépendantes arrivent à la même conclusion.",
        bg=C_BLUE_L, text_color=C_BLUE_D
    )
    add_image_safe(doc, FIG / "13_interpretability_comparison.png", Inches(5.5))
    doc.add_paragraph("Figure 14 — Comparaison des 3 méthodes (normalisées)").alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_heading(doc, "6.5 Recommandations marketing (SHAP → Décision)", level=2, color=C_BLUE_M)
    reco_items = [
        "PRIORISER les investissements TV — correlation 0.999, 95.69% d'importance SHAP. Levier numéro 1.",
        "OPTIMISER la Radio en second plan — impact réel mais marginal par rapport à TV.",
        "NE PAS SURINVESTIR en Social Media — impact faible (0.47% SHAP), rendement décroissant.",
        "NE PAS SURPAYER pour des Mega influenceurs — impact de 0.13% seulement selon SHAP. Aucune justification ROI.",
    ]
    for item in reco_items:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(item)
        run.font.size = Pt(11)
        run.font.name = "Calibri"

    doc.add_page_break()

    # ─────────────────────────────────────────────────────────────────────────
    # 7. DASHBOARD
    # ─────────────────────────────────────────────────────────────────────────
    add_heading(doc, "7. Dashboard décisionnel (Streamlit)", level=1)
    doc.add_paragraph(
        "Le dashboard Streamlit (dashboard/app.py) constitue l'interface décisionnelle "
        "destinée au CMO, responsable marketing ou direction financière. Il est indépendant "
        "du code d'entraînement et charge les modèles sérialisés (joblib). "
        "Lancement : streamlit run dashboard/app.py"
    )
    add_metric_table(doc,
        ["Page", "Contenu"],
        [
            ["Accueil & KPI", "4 indicateurs clés, distribution Sales, répartition Influencer, insights CMO"],
            ["Exploration", "Filtres interactifs (Influencer, TV, Sales), scatter + boxplot, matrice corrélation"],
            ["Simulation budgétaire", "Sliders budgets + Influencer → prédiction temps réel + jauge ROI + histogramme positionnel"],
            ["Comparaison des modèles", "Tableau stylé, bar charts RMSE/R², radar chart multi-critères, scatter prédictions vs réelles"],
            ["Explication SHAP", "Importance globale 3 méthodes, waterfall local par scénario saisi"],
        ]
    )
    doc.add_paragraph(
        "Technologies : Python 3.13, Streamlit 1.52, Plotly 5.x, SHAP 0.45, joblib. "
        "Toutes les prédictions passent par le pipeline complet (preprocessing + modèle), "
        "garantissant une cohérence totale avec l'environnement d'entraînement."
    )

    doc.add_page_break()

    # ─────────────────────────────────────────────────────────────────────────
    # 8. API REST
    # ─────────────────────────────────────────────────────────────────────────
    add_heading(doc, "8. API REST (FastAPI)", level=1)
    doc.add_paragraph(
        "L'API REST expose le modèle XGBoost en tant que service d'inférence. "
        "Elle reprend le pipeline complet sérialisé (preprocessing inclus) pour garantir "
        "une cohérence parfaite avec l'entraînement. "
        "Lancement : uvicorn api.main:app --port 8000\n"
        "Documentation Swagger auto : http://localhost:8000/docs"
    )
    add_metric_table(doc,
        ["Endpoint", "Méthode", "Description"],
        [
            ["/health", "GET", "Statut du service + confirmation que le modèle est chargé"],
            ["/predict", "POST", "Prédiction Sales (M$) + ROI estimé à partir de TV, Radio, SM, Influencer"],
            ["/model-info", "GET", "Métriques du modèle en production (RMSE, R², MAPE, features)"],
        ]
    )
    doc.add_paragraph(
        "Gestion des erreurs : validation Pydantic sur toutes les entrées (valeurs hors plage, "
        "Influencer invalide → HTTP 422 avec message explicite). Le modèle est chargé une seule "
        "fois au démarrage (pas de rechargement à chaque requête)."
    )

    doc.add_page_break()

    # ─────────────────────────────────────────────────────────────────────────
    # 9. CONCLUSION
    # ─────────────────────────────────────────────────────────────────────────
    add_heading(doc, "9. Conclusion et recommandations", level=1)
    add_heading(doc, "9.1 Bilan technique", level=2, color=C_BLUE_M)
    doc.add_paragraph(
        "Ce projet démontre la faisabilité d'un système intelligent complet d'aide à la décision "
        "pour l'optimisation marketing. Le pipeline end-to-end — de l'EDA brute à l'API déployable "
        "— est opérationnel, reproductible (random_state=42 partout) et exempt de data leakage.\n\n"
        "Le modèle XGBoost retenu (RMSE=3.538, R²=0.9985, MAPE=1.84%) est suffisamment précis "
        "pour être utilisé en production pour la planification budgétaire. L'analyse SHAP révèle "
        "une dépendance quasi exclusive aux budgets TV, ce qui simplifie considérablement "
        "les décisions d'allocation."
    )
    add_heading(doc, "9.2 Résultat inattendu — Le Deep Learning non supérieur", level=2, color=C_BLUE_M)
    add_colored_box(doc,
        "Le MLP (Deep Learning) est le modèle le moins performant (RMSE=5.983), inférieur même "
        "à la Régression Linéaire (RMSE=5.879). Ce résultat pédagogiquement important illustre "
        "un principe fondamental : le Deep Learning n'est pas universellement supérieur. "
        "Sur un dataset tabulaire de ~4 500 lignes avec des relations essentiellement linéaires, "
        "les modèles basés sur les arbres sont généralement plus efficaces.",
        bg=RGBColor(0xfe, 0xf3, 0xc7), text_color=RGBColor(0x78, 0x35, 0x00)
    )
    add_heading(doc, "9.3 Perspectives d'amélioration", level=2, color=C_BLUE_M)
    for item in [
        "Enrichir le dataset avec des données temporelles pour capturer les effets de saisonnalité",
        "Implémenter un monitoring de la dérive des données en production (data drift detection)",
        "Tester LightGBM et CatBoost comme alternatives à XGBoost",
        "Déployer l'API sur un cloud provider (AWS, Azure, GCP) avec Docker + CI/CD",
        "Ajouter une tâche de classification (High/Medium/Low Performance) en complément",
    ]:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(item)
        run.font.size = Pt(11)
        run.font.name = "Calibri"

    doc.add_page_break()

    add_heading(doc, "Références", level=1)
    for ref in [
        "Kaggle Dataset : harrimansaragih/dummy-advertising-and-sales-data",
        "Scikit-learn Documentation — https://scikit-learn.org",
        "XGBoost Documentation — https://xgboost.readthedocs.io",
        "SHAP Documentation — https://shap.readthedocs.io",
        "Streamlit Documentation — https://docs.streamlit.io",
        "FastAPI Documentation — https://fastapi.tiangolo.com",
        "Lundberg & Lee (2017) — A Unified Approach to Interpreting Model Predictions (SHAP)",
    ]:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(ref)
        run.font.size = Pt(10)
        run.font.name = "Calibri"
        run.font.color.rgb = C_MUTED

    out_path = REPORTS / "Rapport_ROI_Marketing.docx"
    doc.save(out_path)
    print(f"[OK] Rapport sauvegarde : {out_path}")
    return out_path


# =============================================================================
# PRÉSENTATION POWERPOINT
# =============================================================================

def rgb(r, g, b):
    return PptxRGB(r, g, b)


PBLUE_D = rgb(0x1e, 0x3a, 0x5f)
PBLUE_M = rgb(0x25, 0x63, 0xeb)
PBLUE_L = rgb(0xdb, 0xea, 0xfe)
PGREEN  = rgb(0x10, 0xb9, 0x81)
PAMBER  = rgb(0xf5, 0x9e, 0x0b)
PRED    = rgb(0xe7, 0x4c, 0x3c)
PWHITE  = rgb(0xff, 0xff, 0xff)
PGRAY   = rgb(0xf1, 0xf5, 0xf9)
PDARK   = rgb(0x1e, 0x29, 0x3b)
PMUTED  = rgb(0x64, 0x74, 0x8b)

W = PInches(13.33)
H = PInches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        PInches(left), PInches(top), PInches(width), PInches(height)
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 bold=False, italic=False, size=18, color=None,
                 align=None, wrap=True, font_name="Calibri"):
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(
        PInches(left), PInches(top), PInches(width), PInches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    para = tf.paragraphs[0]
    if align == "center":
        para.alignment = PP_ALIGN.CENTER
    elif align == "right":
        para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = PPt(size)
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox


def add_image_to_slide(slide, path, left, top, width):
    p = Path(path)
    if p.exists():
        slide.shapes.add_picture(str(p), PInches(left), PInches(top), PInches(width))


def ppt_divider(slide, y=1.05, color=None):
    shape = slide.shapes.add_shape(1, PInches(0.7), PInches(y), PInches(11.93), PInches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color or PBLUE_L
    shape.line.fill.background()


def build_presentation():
    prs = new_prs()

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 1 — TITRE
    # ─────────────────────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    # Fond bleu foncé
    add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=PBLUE_D)
    # Accent décoratif
    add_rect(sl, 0, 5.8, 13.33, 1.7, fill_rgb=rgb(0x0f, 0x23, 0x47))
    add_rect(sl, 0.7, 2.8, 0.08, 2.0, fill_rgb=PBLUE_M)
    # Texte
    add_text_box(sl, "EFREI Paris  |  DATA Engineering & AI  |  M1  |  2025–2026",
                 1.0, 0.55, 11.5, 0.5, size=12, color=rgb(0x94, 0xa3, 0xb8))
    add_text_box(sl, "Système Intelligent Multi-Modèles",
                 1.0, 1.2, 11.0, 1.0, bold=True, size=30, color=PWHITE)
    add_text_box(sl, "Optimisation du Retour sur Investissement Marketing",
                 1.0, 2.15, 11.0, 0.9, bold=True, size=22, color=PBLUE_M)
    add_text_box(sl, "Prédiction des ventes multicanal · Machine Learning & Deep Learning · Dashboard · API REST",
                 1.0, 3.1, 11.0, 0.7, size=13, color=rgb(0xdb, 0xea, 0xfe))
    add_text_box(sl, "Salah Eddine Zbiri  ·  Encadrant : Sarah Malaeb  ·  Juin 2026",
                 1.0, 6.2, 11.0, 0.6, size=12, color=rgb(0x94, 0xa3, 0xb8))

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 2 — AGENDA
    # ─────────────────────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=PGRAY)
    add_rect(sl, 0, 0, 13.33, 1.2, fill_rgb=PBLUE_D)
    add_text_box(sl, "Agenda", 0.5, 0.25, 12.0, 0.8, bold=True, size=26, color=PWHITE)
    items = [
        ("01", "Contexte & Problématique"),
        ("02", "Dataset & Analyse Exploratoire (EDA)"),
        ("03", "Pipeline Méthodologique"),
        ("04", "Comparaison des 4 Modèles"),
        ("05", "Résultats & Sélection du Modèle Final"),
        ("06", "Interprétabilité SHAP"),
        ("07", "Dashboard Décisionnel"),
        ("08", "API REST"),
        ("09", "Conclusion & Recommandations"),
    ]
    cols = [(0.6, 4.8), (6.9, 4.8)]
    for i, (num, title) in enumerate(items):
        col_x, col_w = cols[i % 2]
        row = i // 2
        y = 1.4 + row * 1.0
        add_rect(sl, col_x, y, 0.55, 0.55, fill_rgb=PBLUE_M)
        add_text_box(sl, num, col_x, y, 0.55, 0.55, bold=True, size=13, color=PWHITE, align="center")
        add_text_box(sl, title, col_x + 0.65, y + 0.05, col_w - 0.7, 0.5, size=13, color=PDARK)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 3 — CONTEXTE
    # ─────────────────────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=PGRAY)
    add_rect(sl, 0, 0, 13.33, 1.2, fill_rgb=PBLUE_D)
    add_text_box(sl, "01 — Contexte & Problématique", 0.5, 0.25, 12.0, 0.8, bold=True, size=22, color=PWHITE)
    ppt_divider(sl, y=1.25, color=PBLUE_M)
    add_text_box(sl,
        "Les entreprises investissent des millions en publicité sans outil fiable pour estimer\n"
        "leur retour sur investissement AVANT l'exécution des campagnes.",
        0.7, 1.35, 11.9, 0.9, size=14, color=PDARK)
    challenges = [
        ("TV", "Budget le plus impactant mais le plus coûteux", PRED),
        ("Radio", "Impact secondaire, complémentaire à TV", PAMBER),
        ("Social Media", "Impact modéré, difficile à mesurer", PBLUE_M),
        ("Influencer", "Coût élevé, impact réel incertain", PGREEN),
    ]
    for i, (canal, desc, col) in enumerate(challenges):
        x = 0.7 + i * 3.0
        add_rect(sl, x, 2.4, 2.7, 1.5, fill_rgb=PWHITE)
        add_rect(sl, x, 2.4, 2.7, 0.08, fill_rgb=col)
        add_text_box(sl, canal, x + 0.1, 2.55, 2.5, 0.5, bold=True, size=16, color=col)
        add_text_box(sl, desc, x + 0.1, 3.0, 2.5, 0.8, size=11, color=PMUTED)
    add_rect(sl, 0.7, 4.2, 11.93, 1.4, fill_rgb=PBLUE_D)
    add_text_box(sl,
        "Question centrale : Comment prédire le volume de ventes généré par une combinaison\n"
        "de budgets marketing et identifier les canaux qui maximisent réellement le ROI ?",
        0.9, 4.3, 11.5, 1.1, bold=True, size=15, color=PWHITE)
    add_text_box(sl,
        "Tâche choisie : Régression (prédiction de Sales en M$) · Variable cible : Sales",
        0.7, 5.8, 11.9, 0.5, size=12, color=PMUTED, italic=True)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 4 — DATASET & EDA
    # ─────────────────────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=PGRAY)
    add_rect(sl, 0, 0, 13.33, 1.2, fill_rgb=PBLUE_D)
    add_text_box(sl, "02 — Dataset & Analyse Exploratoire", 0.5, 0.25, 12.0, 0.8, bold=True, size=22, color=PWHITE)
    # Stats clés
    kpis = [
        ("~4 500", "Campagnes", PBLUE_M),
        ("4", "Features", PGREEN),
        ("1", "Variable cible", PAMBER),
        ("0.999", "Corr. TV→Sales", PRED),
    ]
    for i, (val, label, col) in enumerate(kpis):
        x = 0.5 + i * 3.1
        add_rect(sl, x, 1.4, 2.8, 1.5, fill_rgb=PWHITE)
        add_rect(sl, x, 1.4, 2.8, 0.07, fill_rgb=col)
        add_text_box(sl, val, x + 0.1, 1.55, 2.6, 0.7, bold=True, size=28, color=col)
        add_text_box(sl, label, x + 0.1, 2.2, 2.6, 0.4, size=11, color=PMUTED)
    add_image_to_slide(sl, FIG / "05_budgets_vs_sales.png", 0.4, 3.1, 7.8)
    add_image_to_slide(sl, FIG / "07_correlation_matrix.png", 8.5, 3.0, 4.4)
    add_text_box(sl, "Budget TV vs Sales (corr. = 0.999)", 0.4, 6.5, 7.8, 0.4, size=10, color=PMUTED, italic=True)
    add_text_box(sl, "Matrice de corrélation", 8.5, 6.5, 4.4, 0.4, size=10, color=PMUTED, italic=True)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 5 — MÉTHODOLOGIE
    # ─────────────────────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=PGRAY)
    add_rect(sl, 0, 0, 13.33, 1.2, fill_rgb=PBLUE_D)
    add_text_box(sl, "03 — Pipeline Méthodologique", 0.5, 0.25, 12.0, 0.8, bold=True, size=22, color=PWHITE)
    steps = [
        ("EDA", "Distributions, corrélations, outliers, insights métier", PBLUE_M),
        ("Preprocessing", "Imputation médiane, StandardScaler,\nOneHotEncoder · FIT sur train only", PGREEN),
        ("Modélisation", "4 modèles · GridSearchCV\n5-fold KFold · Pipeline sklearn", PAMBER),
        ("Évaluation", "MAE, RMSE, R², MAPE\nTrain vs Test · Overfitting check", PRED),
        ("Interprétabilité", "Feature Importance\nPermutation · SHAP", rgb(0x53, 0x34, 0x83)),
    ]
    arrow_x_start = 0.55
    box_w = 2.25
    gap = 0.22
    for i, (title, desc, col) in enumerate(steps):
        x = arrow_x_start + i * (box_w + gap)
        add_rect(sl, x, 1.6, box_w, 2.5, fill_rgb=PWHITE)
        add_rect(sl, x, 1.6, box_w, 0.08, fill_rgb=col)
        add_rect(sl, x, 1.6, box_w, 0.5, fill_rgb=col)
        add_text_box(sl, f"0{i+1}", x + 0.07, 1.62, 0.4, 0.4, bold=True, size=16, color=PWHITE)
        add_text_box(sl, title, x + 0.07, 2.18, box_w - 0.1, 0.5, bold=True, size=13, color=PDARK)
        add_text_box(sl, desc, x + 0.07, 2.68, box_w - 0.1, 1.3, size=10, color=PMUTED)
        if i < len(steps) - 1:
            add_rect(sl, x + box_w + 0.03, 2.65, 0.16, 0.16, fill_rgb=PBLUE_M)

    add_rect(sl, 0.55, 4.45, 12.23, 1.0, fill_rgb=PBLUE_D)
    add_text_box(sl,
        "Principe fondamental ANTI DATA LEAKAGE : le preprocessor est FIT uniquement sur le train set "
        "et appliqué en transform sur le test — les pipelines sklearn garantissent ce principe même "
        "pendant la cross-validation (refit à chaque fold).",
        0.75, 4.5, 11.8, 0.9, size=12, color=PWHITE)
    add_text_box(sl, "Livrables : 5 notebooks · src/ réutilisable · dashboard Streamlit · API FastAPI",
                 0.55, 5.65, 12.0, 0.5, size=11, color=PMUTED, italic=True)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 6 — 4 MODÈLES
    # ─────────────────────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=PGRAY)
    add_rect(sl, 0, 0, 13.33, 1.2, fill_rgb=PBLUE_D)
    add_text_box(sl, "04 — Les 4 Modèles Comparés", 0.5, 0.25, 12.0, 0.8, bold=True, size=22, color=PWHITE)
    models_ppt = [
        ("Régression\nLinéaire", "Baseline · Interprétable\nAucun tuning", f"CV RMSE: {train_df[train_df['Model']=='Linear Regression']['CV_RMSE'].values[0]:.3f}", PBLUE_M),
        ("Random\nForest", "Ensemble · 200 arbres\nCapture non-linéarités", f"CV RMSE: {train_df[train_df['Model']=='Random Forest']['CV_RMSE'].values[0]:.3f}", PGREEN),
        ("XGBoost", "Gradient Boosting\nÉtat de l'art tabulaire\n★ RETENU", f"CV RMSE: {train_df[train_df['Model']=='XGBoost']['CV_RMSE'].values[0]:.3f}", PAMBER),
        ("MLP\n(Deep Learning)", "Réseau de neurones\nEarly stopping + L2\nArchitecture (50,25)", f"CV RMSE: {train_df[train_df['Model']=='MLP (Deep Learning)']['CV_RMSE'].values[0]:.3f}", PRED),
    ]
    for i, (name, desc, cv, col) in enumerate(models_ppt):
        x = 0.55 + i * 3.15
        add_rect(sl, x, 1.45, 2.85, 3.7, fill_rgb=PWHITE)
        add_rect(sl, x, 1.45, 2.85, 0.07, fill_rgb=col)
        add_text_box(sl, name, x + 0.12, 1.6, 2.6, 0.8, bold=True, size=15, color=col)
        add_text_box(sl, desc, x + 0.12, 2.4, 2.6, 1.2, size=11, color=PMUTED)
        add_text_box(sl, cv, x + 0.12, 3.55, 2.6, 0.45, bold=True, size=13, color=PDARK)
        if i == 2:  # XGBoost — highlight
            add_rect(sl, x, 5.25, 2.85, 0.55, fill_rgb=PGREEN)
            add_text_box(sl, "MODELE FINAL", x + 0.12, 5.3, 2.6, 0.4, bold=True, size=12, color=PWHITE, align="center")
    add_text_box(sl,
        "GridSearchCV 5-fold KFold · Métrique de sélection : RMSE · Pipeline complet (preprocessor + estimator)",
        0.55, 5.95, 12.2, 0.5, size=11, color=PMUTED, italic=True)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 7 — RÉSULTATS
    # ─────────────────────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=PGRAY)
    add_rect(sl, 0, 0, 13.33, 1.2, fill_rgb=PBLUE_D)
    add_text_box(sl, "05 — Résultats — Test Set (données jamais vues)", 0.5, 0.25, 12.0, 0.8, bold=True, size=22, color=PWHITE)
    # Tableau
    headers_t = ["Modèle", "MAE", "RMSE", "R²", "MAPE", "Verdict"]
    rows_t = []
    for _, row in eval_df.sort_values("RMSE_test").iterrows():
        rows_t.append([row["Model"], f"{row['MAE_test']:.3f}", f"{row['RMSE_test']:.3f}",
                       f"{row['R2_test']:.4f}", f"{row['MAPE_test_%']:.2f}%", row["Verdict"]])
    col_widths = [3.0, 1.5, 1.5, 1.5, 1.5, 1.5]
    col_colors = [PBLUE_M] * 6
    y_table = 1.35
    # Header row
    x_cur = 0.55
    for j, (hdr, cw) in enumerate(zip(headers_t, col_widths)):
        add_rect(sl, x_cur, y_table, cw, 0.45, fill_rgb=PBLUE_D)
        add_text_box(sl, hdr, x_cur + 0.05, y_table + 0.05, cw - 0.1, 0.35,
                     bold=True, size=11, color=PWHITE, align="center")
        x_cur += cw
    for ri, row_data in enumerate(rows_t):
        y_r = y_table + 0.45 + ri * 0.45
        bg = PWHITE if ri % 2 == 0 else PGRAY
        if ri == 0:  # XGBoost
            bg = rgb(0xd1, 0xfa, 0xe5)
        x_cur = 0.55
        for j, (val, cw) in enumerate(zip(row_data, col_widths)):
            add_rect(sl, x_cur, y_r, cw, 0.45, fill_rgb=bg)
            col_txt = PGREEN if ri == 0 else PDARK
            add_text_box(sl, val, x_cur + 0.05, y_r + 0.05, cw - 0.1, 0.35,
                         bold=(ri == 0), size=11, color=col_txt, align="center")
            x_cur += cw
    add_image_to_slide(sl, FIG / "07_metrics_comparison.png", 0.4, 4.35, 8.0)
    add_rect(sl, 8.7, 4.35, 4.3, 2.7, fill_rgb=PBLUE_D)
    add_text_box(sl, "XGBoost retenu", 8.8, 4.4, 4.1, 0.5, bold=True, size=15, color=PGREEN)
    best = eval_df[eval_df["Model"] == "XGBoost"].iloc[0]
    for li, (label, val) in enumerate([("RMSE", f"{best['RMSE_test']:.3f}"), ("R²", f"{best['R2_test']:.4f}"),
                                         ("MAPE", f"{best['MAPE_test_%']:.2f}%"), ("Overfitting", "Aucun ✓")]):
        add_text_box(sl, f"{label}  :  {val}", 8.8, 5.0 + li * 0.5, 4.1, 0.45,
                     size=13, color=PWHITE if li > 0 else PGREEN, bold=(li == 0))

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 8 — SHAP
    # ─────────────────────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=PGRAY)
    add_rect(sl, 0, 0, 13.33, 1.2, fill_rgb=PBLUE_D)
    add_text_box(sl, "06 — Interprétabilité SHAP — 3 méthodes", 0.5, 0.25, 12.0, 0.8, bold=True, size=22, color=PWHITE)
    add_image_to_slide(sl, FIG / "10_shap_importance_bar.png", 0.4, 1.3, 6.2)
    add_image_to_slide(sl, FIG / "13_interpretability_comparison.png", 6.8, 1.3, 6.1)
    add_text_box(sl, "SHAP Importance globale (mean |SHAP value|)", 0.4, 5.3, 6.2, 0.4, size=9, color=PMUTED, italic=True)
    add_text_box(sl, "Comparaison 3 méthodes (normalisées à 100%)", 6.8, 5.3, 6.1, 0.4, size=9, color=PMUTED, italic=True)
    insights_shap = [
        ("TV", "95.69%", PRED),
        ("Radio", "3.71%", PAMBER),
        ("Social Media", "0.47%", PBLUE_M),
        ("Influencer", "0.13%", PGREEN),
    ]
    for i, (feat, pct, col) in enumerate(insights_shap):
        x = 0.5 + i * 3.1
        add_rect(sl, x, 5.8, 2.85, 0.9, fill_rgb=PWHITE)
        add_rect(sl, x, 5.8, 2.85, 0.07, fill_rgb=col)
        add_text_box(sl, feat, x + 0.1, 5.9, 1.5, 0.4, bold=True, size=13, color=PDARK)
        add_text_box(sl, pct, x + 1.5, 5.9, 1.3, 0.4, bold=True, size=14, color=col, align="right")
        add_text_box(sl, "SHAP importance", x + 0.1, 6.28, 2.6, 0.35, size=9, color=PMUTED)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 9 — DASHBOARD
    # ─────────────────────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=PGRAY)
    add_rect(sl, 0, 0, 13.33, 1.2, fill_rgb=PBLUE_D)
    add_text_box(sl, "07 — Dashboard Décisionnel (Streamlit)", 0.5, 0.25, 12.0, 0.8, bold=True, size=22, color=PWHITE)
    pages_dash = [
        ("Accueil & KPI", "KPI executives, distribution Sales, insights CMO"),
        ("Exploration", "Filtres interactifs, scatter, boxplot, heatmap"),
        ("Simulation", "Sliders budgets → prédiction + jauge ROI temps réel"),
        ("Comparaison", "Tableau stylé, radar chart, prédictions vs réelles"),
        ("SHAP Local", "Waterfall SHAP — expliquer UNE prédiction"),
    ]
    for i, (title, desc) in enumerate(pages_dash):
        x = 0.5 + (i % 3) * 4.1
        y = 1.5 + (i // 3) * 2.3
        col = [PBLUE_M, PGREEN, PAMBER, PRED, rgb(0x53, 0x34, 0x83)][i]
        add_rect(sl, x, y, 3.7, 1.9, fill_rgb=PWHITE)
        add_rect(sl, x, y, 3.7, 0.07, fill_rgb=col)
        add_rect(sl, x, y, 0.55, 1.9, fill_rgb=col)
        add_text_box(sl, str(i + 1), x + 0.07, y + 0.65, 0.4, 0.55,
                     bold=True, size=18, color=PWHITE, align="center")
        add_text_box(sl, title, x + 0.65, y + 0.15, 2.95, 0.55, bold=True, size=12, color=PDARK)
        add_text_box(sl, desc, x + 0.65, y + 0.7, 2.95, 1.0, size=10, color=PMUTED)
    add_rect(sl, 0.5, 6.05, 12.3, 0.65, fill_rgb=PBLUE_D)
    add_text_box(sl,
        "http://localhost:8501  ·  streamlit run dashboard/app.py  ·  Modèle : XGBoost R²=0.9985",
        0.7, 6.1, 11.9, 0.5, size=12, color=PWHITE)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 10 — API REST
    # ─────────────────────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=PGRAY)
    add_rect(sl, 0, 0, 13.33, 1.2, fill_rgb=PBLUE_D)
    add_text_box(sl, "08 — API REST (FastAPI) — Industrialisation", 0.5, 0.25, 12.0, 0.8, bold=True, size=22, color=PWHITE)
    endpoints = [
        ("GET", "/health", "Statut du service + confirmation modèle chargé", PGREEN),
        ("POST", "/predict", "Prédiction Sales (M$) + ROI estimé + validation entrées", PAMBER),
        ("GET", "/model-info", "Métriques production : RMSE, R², MAPE, features", PBLUE_M),
    ]
    for i, (method, path, desc, col) in enumerate(endpoints):
        y = 1.5 + i * 1.3
        add_rect(sl, 0.5, y, 1.1, 0.55, fill_rgb=col)
        add_text_box(sl, method, 0.5, y + 0.05, 1.1, 0.45, bold=True, size=13, color=PWHITE, align="center")
        add_rect(sl, 1.7, y, 2.8, 0.55, fill_rgb=rgb(0x1e, 0x29, 0x3b))
        add_text_box(sl, path, 1.75, y + 0.05, 2.7, 0.45, bold=True, size=13, color=PWHITE, font_name="Courier New")
        add_text_box(sl, desc, 4.65, y + 0.08, 8.1, 0.45, size=12, color=PDARK)
    # Exemple JSON
    add_rect(sl, 0.5, 5.5, 5.5, 1.6, fill_rgb=rgb(0x1e, 0x29, 0x3b))
    add_text_box(sl, "Request", 0.5, 5.25, 2.5, 0.35, bold=True, size=11, color=PMUTED)
    req_json = '{\n  "TV": 230.1,\n  "Radio": 37.8,\n  "Social Media": 69.2,\n  "Influencer": "Macro"\n}'
    add_text_box(sl, req_json, 0.65, 5.55, 5.2, 1.5, size=10, color=PWHITE, font_name="Courier New")
    add_rect(sl, 6.3, 5.5, 6.5, 1.6, fill_rgb=rgb(0x06, 0x5f, 0x46))
    add_text_box(sl, "Response", 6.3, 5.25, 3.0, 0.35, bold=True, size=11, color=PMUTED)
    resp_json = '{\n  "predicted_sales_M": 355.43,\n  "roi_estimated": 1.054,\n  "model_used": "XGBoost"\n}'
    add_text_box(sl, resp_json, 6.45, 5.55, 6.2, 1.5, size=10, color=PWHITE, font_name="Courier New")

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 11 — CONCLUSION
    # ─────────────────────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=PGRAY)
    add_rect(sl, 0, 0, 13.33, 1.2, fill_rgb=PBLUE_D)
    add_text_box(sl, "09 — Conclusion & Recommandations", 0.5, 0.25, 12.0, 0.8, bold=True, size=22, color=PWHITE)
    # Bilan
    add_rect(sl, 0.5, 1.35, 12.3, 0.6, fill_rgb=rgb(0xd1, 0xfa, 0xe5))
    add_text_box(sl,
        "XGBoost : RMSE=3.538 · R²=0.9985 · MAPE=1.84% · Zéro overfitting — Précision suffisante pour planification budgétaire",
        0.65, 1.38, 11.9, 0.5, bold=True, size=13, color=rgb(0x06, 0x5f, 0x46))
    recos = [
        ("Prioriser le budget TV", "Corrélation 0.999 et 95.69% d'importance SHAP. ROI le plus élevé.", PRED),
        ("Optimiser Radio en second", "Impact réel mais marginal (3.71% SHAP). Complémentaire à TV.", PAMBER),
        ("Limiter Social Media", "Impact faible (0.47% SHAP). Rendement décroissant au-delà d'un seuil.", PBLUE_M),
        ("Ne pas surpayer l'Influencer", "0.13% d'importance. Mega ≈ Nano en termes d'impact sur les ventes.", PGREEN),
    ]
    for i, (title, desc, col) in enumerate(recos):
        x = 0.5 + (i % 2) * 6.2
        y = 2.2 + (i // 2) * 1.8
        add_rect(sl, x, y, 5.9, 1.5, fill_rgb=PWHITE)
        add_rect(sl, x, y, 5.9, 0.07, fill_rgb=col)
        add_text_box(sl, title, x + 0.12, y + 0.15, 5.65, 0.5, bold=True, size=13, color=col)
        add_text_box(sl, desc, x + 0.12, y + 0.62, 5.65, 0.75, size=11, color=PMUTED)
    key_insight = (
        "Résultat inattendu : le MLP (Deep Learning) est INFÉRIEUR à la Régression Linéaire.\n"
        "Sur données tabulaires avec relations quasi-linéaires, les modèles à base d'arbres dominent."
    )
    add_rect(sl, 0.5, 5.95, 12.3, 0.85, fill_rgb=rgb(0xfe, 0xf3, 0xc7))
    add_text_box(sl, key_insight, 0.65, 5.98, 11.9, 0.8, size=11, color=rgb(0x78, 0x35, 0x00))

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 12 — MERCI
    # ─────────────────────────────────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=PBLUE_D)
    add_rect(sl, 0, 5.5, 13.33, 2.0, fill_rgb=rgb(0x0f, 0x23, 0x47))
    add_rect(sl, 5.8, 2.0, 0.08, 2.8, fill_rgb=PBLUE_M)
    add_text_box(sl, "Merci pour votre attention", 0.6, 1.8, 12.0, 1.0,
                 bold=True, size=34, color=PWHITE, align="center")
    add_text_box(sl, "Questions ?", 0.6, 2.9, 12.0, 0.8,
                 size=22, color=PBLUE_M, align="center")
    add_text_box(sl,
        "Salah Eddine Zbiri  ·  claudezbiri@gmail.com  ·  EFREI M1 DE 2025/2026",
        0.6, 6.1, 12.0, 0.5, size=12, color=rgb(0x94, 0xa3, 0xb8), align="center")
    add_text_box(sl,
        "Dashboard : localhost:8501  ·  API : localhost:8000/docs  ·  Code : GitHub",
        0.6, 6.65, 12.0, 0.5, size=11, color=rgb(0x64, 0x74, 0x8b), align="center")

    out_path = REPORTS / "Presentation_ROI_Marketing.pptx"
    prs.save(out_path)
    print(f"[OK] Presentation sauvegardee : {out_path}")
    return out_path


# =============================================================================
# MAIN
# =============================================================================

if __name__ == "__main__":
    print("Generation des livrables...")
    print()
    r = build_report()
    print()
    p = build_presentation()
    print()
    print("=" * 60)
    print("LIVRABLES GENERES :")
    print(f"  Rapport      : {r}")
    print(f"  Presentation : {p}")
    print("=" * 60)
